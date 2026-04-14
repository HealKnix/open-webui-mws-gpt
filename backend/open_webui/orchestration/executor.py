from __future__ import annotations

import asyncio
import hashlib
import io
import json
import logging
import os
import re
import subprocess
import sys
import tempfile
import time
import uuid
from pathlib import Path
from urllib.parse import urlparse
from typing import Any, Awaitable, Callable

import tiktoken
from fastapi import HTTPException, Request, status
from starlette.responses import StreamingResponse

from open_webui.models.chats import Chats
from open_webui.models.files import FileForm, Files
from open_webui.models.users import UserModel
from open_webui.orchestration.adapters.mws import create_text_response, transcribe_audio
from open_webui.orchestration.planner import build_plan
from open_webui.orchestration.policy import select_operation_assignment
from open_webui.orchestration.registry import MTS_ROUTER_MODEL_ID, build_registry
from open_webui.retrieval.utils import get_content_from_url, get_sources_from_items
from open_webui.routers.retrieval import search_web as retrieval_search_web
from open_webui.orchestration.schema import Artifact, ExecutionResult, ExecutionTraceEvent, FinalOutput
from open_webui.socket.main import get_event_emitter
from open_webui.storage.provider import Storage
from open_webui.utils.models import get_all_models

log = logging.getLogger(__name__)

StatusEmitter = Callable[[dict[str, Any]], Awaitable[None]]

MODEL_CONTEXT_WINDOW_HINTS: dict[str, int] = {
    'qwen2.5-72b-instruct': 32_768,
    'qwen3-32b': 32_768,
    'glm-4.6-357b': 32_768,
    'llama-3.3-70b-instruct': 32_768,
    'mws-gpt-alpha': 32_768,
}

PRESENTATION_CORE_COMPONENTS: tuple[str, ...] = (
    'структура',
    'ключевые цифры',
    'риски',
    'выводы',
    'источники',
)

PRESENTATION_RENDERER_ENV = 'MTS_ROUTER_PPTX_RENDERER'
PRESENTATION_RENDERER_SKILL_PATH_ENV = 'MTS_ROUTER_PPTX_SKILL_PATH'
PRESENTATION_TITLE_FONT_PT = 26
PRESENTATION_BODY_FONT_PT = 13
PRESENTATION_SOURCES_FONT_PT = 8


def _add_trace(result: ExecutionResult, phase: str, status_name: str, detail: str | None = None, **kwargs):
    result.trace.events.append(
        ExecutionTraceEvent(
            phase=phase,
            status=status_name,
            detail=detail,
            **kwargs,
        )
    )


def _artifact_map(result: ExecutionResult) -> dict[str, Artifact]:
    return {artifact.id: artifact for artifact in [*result.plan.input_artifacts, *result.output_artifacts]}


def _get_text(artifact_map: dict[str, Artifact], artifact_id: str | None) -> str:
    if not artifact_id:
        return ''
    artifact = artifact_map.get(artifact_id)
    return artifact.text if artifact and artifact.text else ''


def _build_context_block(artifact_map: dict[str, Artifact], input_ids: list[str]) -> str:
    blocks = []
    for input_id in input_ids:
        artifact = artifact_map.get(input_id)
        if not artifact:
            continue
        if artifact.type == 'text' and artifact.text:
            blocks.append(f'[text]\n{artifact.text}')
        elif artifact.type == 'structured' and artifact.text:
            blocks.append(f'[structured]\n{artifact.text}')
        elif artifact.type == 'image' and artifact.text:
            blocks.append(f'[image]\n{artifact.text}')
        elif artifact.type == 'audio' and artifact.text:
            blocks.append(f'[audio]\n{artifact.text}')
    return '\n\n'.join(blocks)


def _truncate_error(exc: Exception, max_length: int = 220) -> str:
    text = str(exc).strip()
    if len(text) <= max_length:
        return text
    return f'{text[:max_length].rstrip()}...'


def _is_transient_provider_error(exc: Exception) -> bool:
    text = str(exc or '').lower()
    if not text:
        return False
    transient_markers = (
        'server disconnected',
        'connection error',
        'cannot connect',
        'timed out',
        'timeout',
        'temporarily unavailable',
        'service unavailable',
        '502',
        '503',
        '504',
    )
    return any(marker in text for marker in transient_markers)


def _context_window_hint_for_model(model_id: str) -> int:
    lowered = (model_id or '').strip().lower()
    for key, limit in MODEL_CONTEXT_WINDOW_HINTS.items():
        if key.lower() in lowered:
            return limit
    return 32_768


def _get_user_attr(user: UserModel | Any, attr: str, fallback: str = '') -> str:
    value = None
    if isinstance(user, dict):
        value = user.get(attr)
    else:
        value = getattr(user, attr, None)
    if value is None:
        return fallback
    return str(value)


def _safe_pptx_filename(title: str | None) -> str:
    base = re.sub(r'\s+', ' ', (title or '').strip())
    if not base:
        base = f'presentation-{time.strftime("%Y%m%d-%H%M%S")}'
    base = re.sub(r'[^0-9A-Za-zА-Яа-я._ -]+', '_', base).strip(' ._-')
    if not base:
        base = f'presentation-{uuid.uuid4().hex[:8]}'
    if not base.lower().endswith('.pptx'):
        base = f'{base[:72]}.pptx'
    return base


def _normalize_slide_item(raw_item: Any, index: int) -> dict[str, Any] | None:
    if not isinstance(raw_item, dict):
        return None

    def _clean(value: Any) -> str:
        candidate = str(value or '').strip()
        if candidate.lower() in {'', 'none', 'null', 'nan', 'n/a'}:
            return ''
        return candidate

    title = str(raw_item.get('title') or raw_item.get('name') or f'Слайд {index}').strip()
    text = _clean(raw_item.get('text') or raw_item.get('body'))

    bullets_raw = raw_item.get('list')
    if bullets_raw is None:
        bullets_raw = raw_item.get('bullets')
    if bullets_raw is None:
        bullets_raw = raw_item.get('points')

    bullets: list[str] = []
    if isinstance(bullets_raw, str):
        candidate = bullets_raw.strip()
        if candidate:
            bullets.append(candidate)
    elif isinstance(bullets_raw, list):
        for item in bullets_raw:
            candidate = _clean(item)
            if candidate:
                bullets.append(candidate)

    content_raw = raw_item.get('content')
    content_items: list[str] = []
    if isinstance(content_raw, str):
        candidate = _clean(content_raw)
        if candidate:
            content_items.append(candidate)
    elif isinstance(content_raw, list):
        for item in content_raw:
            if isinstance(item, dict):
                candidate = _clean(
                    item.get('text')
                    or item.get('content')
                    or item.get('fact')
                    or item.get('insight')
                    or item.get('takeaway')
                )
            else:
                candidate = _clean(item)
            if candidate:
                content_items.append(candidate)

    narrative_hints: list[str] = []
    for key in (
        'analysis',
        'insight',
        'takeaway',
        'implication',
        'why_it_matters',
        'commentary',
        'key_message',
    ):
        candidate = _clean(raw_item.get(key))
        if candidate:
            narrative_hints.append(candidate)

    if not text and content_items:
        text = content_items[0]
        if len(content_items) > 1 and not bullets:
            bullets.extend(content_items[1:])
    if not text and narrative_hints:
        text = narrative_hints[0]
    if text and len(text) < 90 and narrative_hints:
        for hint in narrative_hints:
            if hint.lower() not in text.lower():
                text = f'{text} {hint}'.strip()
                break
    if text and len(text) < 90 and len(content_items) > 1:
        for item in content_items[1:]:
            if item.lower() not in text.lower():
                text = f'{text} {item}'.strip()
                if len(text) >= 120:
                    break
    if text and len(text) < 70 and len(bullets) >= 2:
        text = f'{text} Ключевые аспекты: {bullets[0]}; {bullets[1]}.'.strip()
    if not text and bullets:
        text = bullets[0]
        bullets = bullets[1:]

    bullets = bullets[:12]

    sources_raw = raw_item.get('sources') or []
    sources: list[str] = []
    if isinstance(sources_raw, str):
        candidate = sources_raw.strip()
        if candidate:
            sources.append(candidate)
    elif isinstance(sources_raw, list):
        for item in sources_raw:
            candidate = str(item or '').strip()
            if candidate:
                sources.append(candidate)
    seen_sources: set[str] = set()
    dedup_sources: list[str] = []
    for source in sources:
        lowered = source.lower()
        if lowered in seen_sources:
            continue
        seen_sources.add(lowered)
        dedup_sources.append(source)

    return {
        'title': title or f'Слайд {index}',
        'text': text,
        'list': bullets,
        'speaker_notes': str(raw_item.get('speaker_notes') or raw_item.get('notes') or '').strip(),
        'sources': dedup_sources[:16],
    }


def _extract_presentation_deck(text: str) -> list[dict[str, Any]]:
    if not text:
        return []

    candidates: list[str] = []
    stripped = text.strip()
    if stripped.startswith('[') or stripped.startswith('{'):
        candidates.append(stripped)

    code_blocks = re.findall(r'```(?:json)?\s*([\s\S]*?)```', text, flags=re.IGNORECASE)
    candidates.extend(block.strip() for block in code_blocks if block and block.strip())

    bracket_match = re.search(r'(\[\s*\{[\s\S]*\}\s*\])', text)
    if bracket_match:
        candidates.append(bracket_match.group(1).strip())

    for candidate in candidates:
        try:
            parsed = json.loads(candidate)
        except Exception:
            continue

        slides_raw = None
        if isinstance(parsed, list):
            slides_raw = parsed
        elif isinstance(parsed, dict):
            for key in ('slides', 'deck', 'presentation', 'data'):
                if isinstance(parsed.get(key), list):
                    slides_raw = parsed.get(key)
                    break

        if not isinstance(slides_raw, list):
            continue

        normalized: list[dict[str, Any]] = []
        for idx, item in enumerate(slides_raw, start=1):
            slide = _normalize_slide_item(item, idx)
            if slide is not None:
                normalized.append(slide)

        if normalized:
            return normalized

    return []


def _normalize_deck_for_render(deck: list[dict[str, Any]]) -> list[dict[str, Any]]:
    normalized: list[dict[str, Any]] = []
    for index, item in enumerate(deck or [], start=1):
        slide = _normalize_slide_item(item, index)
        if slide is not None:
            normalized.append(slide)
    return normalized


def _extract_presentation_deck_from_artifacts(
    artifacts: list[Artifact],
    preferred_ids: list[str] | None = None,
    *,
    min_slides: int = 3,
) -> list[dict[str, Any]]:
    artifact_map = {artifact.id: artifact for artifact in artifacts or []}
    ordered_ids: list[str] = []
    seen_ids: set[str] = set()

    for artifact_id in preferred_ids or []:
        if artifact_id in artifact_map and artifact_id not in seen_ids:
            ordered_ids.append(artifact_id)
            seen_ids.add(artifact_id)

    for artifact in artifacts or []:
        if artifact.id not in seen_ids:
            ordered_ids.append(artifact.id)
            seen_ids.add(artifact.id)

    for artifact_id in ordered_ids:
        artifact = artifact_map.get(artifact_id)
        if not artifact or not isinstance(artifact.text, str) or not artifact.text.strip():
            continue
        deck = _extract_presentation_deck(artifact.text)
        if len(deck) >= min_slides:
            return deck

    return []


def _is_presentation_refusal_text(text: str) -> bool:
    lowered = (text or '').strip().lower()
    if not lowered:
        return False

    refusal_markers = (
        'не могу сформировать презентацию',
        'не могу сделать презентацию',
        'не могу подготовить презентацию',
        'не могу',
        'выходит за рамки моих возможностей',
        'требует доступа к интернету',
        'необходимы конкретные данные',
        'финальный deck не может быть сформирован',
        'final deck cannot be formed',
        'cannot generate presentation',
    )
    return any(marker in lowered for marker in refusal_markers)


def _build_presentation_recovery_text(deck: list[dict[str, Any]]) -> str:
    unique_sources: list[str] = []
    seen_sources: set[str] = set()
    for slide in deck:
        for source in slide.get('sources') or []:
            candidate = str(source or '').strip()
            lowered = candidate.lower()
            if not candidate or lowered in seen_sources:
                continue
            seen_sources.add(lowered)
            unique_sources.append(candidate)
            if len(unique_sources) >= 12:
                break
        if len(unique_sources) >= 12:
            break

    source_lines = '\n'.join(f'- {item}' for item in unique_sources) if unique_sources else '- Источники внутри deck JSON'

    return (
        '## Краткий вывод\n'
        f'Презентация сформирована на основе собранного контекста ({len(deck)} слайдов).\n\n'
        '## Рекомендации по подаче\n'
        'Используйте слайды как основу, при необходимости уточните цифры под целевую аудиторию.\n\n'
        '## Источники\n'
        f'{source_lines}\n'
    )


def _build_presentation_pptx_bytes_python_pptx(deck: list[dict[str, Any]]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    def _apply_font_size(text_frame, size_pt: int) -> None:
        for paragraph in text_frame.paragraphs:
            try:
                paragraph.font.size = Pt(size_pt)
            except Exception:
                pass
            for run in paragraph.runs:
                try:
                    run.font.size = Pt(size_pt)
                except Exception:
                    pass

    prs = Presentation()
    for index, slide_item in enumerate(deck, start=1):
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        title_text = str(slide_item.get('title') or f'Слайд {index}')
        content_text = str(slide_item.get('text') or '').strip()
        bullet_items = [str(item).strip() for item in (slide_item.get('list') or []) if str(item).strip()]
        speaker_notes = str(slide_item.get('speaker_notes') or '').strip()
        source_items = [str(item).strip() for item in (slide_item.get('sources') or []) if str(item).strip()]

        if slide.shapes.title:
            slide.shapes.title.text = title_text
            try:
                _apply_font_size(slide.shapes.title.text_frame, PRESENTATION_TITLE_FONT_PT)
            except Exception:
                pass

        text_frame = None
        if len(slide.placeholders) > 1:
            try:
                body_placeholder = slide.placeholders[1]
                if hasattr(body_placeholder, 'text_frame'):
                    text_frame = body_placeholder.text_frame
            except Exception:
                text_frame = None

        if text_frame is None:
            textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.6))
            text_frame = textbox.text_frame

        text_frame.clear()
        if content_text:
            text_frame.text = content_text
        elif bullet_items:
            text_frame.text = bullet_items[0]
            bullet_items = bullet_items[1:]
        else:
            text_frame.text = ''

        for bullet in bullet_items:
            paragraph = text_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0

        _apply_font_size(text_frame, PRESENTATION_BODY_FONT_PT)

        if source_items:
            sources_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.6), Inches(11.8), Inches(0.7))
            sources_tf = sources_box.text_frame
            sources_tf.clear()
            source_preview = '; '.join(source_items[:4])
            sources_tf.text = f'Источники: {source_preview}'
            _apply_font_size(sources_tf, PRESENTATION_SOURCES_FONT_PT)

        if speaker_notes or source_items:
            notes_block: list[str] = []
            if speaker_notes:
                notes_block.append(speaker_notes)
            if source_items:
                notes_block.append('Источники:\n' + '\n'.join(f'- {item}' for item in source_items))
            try:
                slide.notes_slide.notes_text_frame.text = '\n\n'.join(notes_block).strip()
            except Exception:
                pass

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _resolve_pptx_skill_root() -> Path | None:
    raw_env_path = (os.getenv(PRESENTATION_RENDERER_SKILL_PATH_ENV) or '').strip()
    candidates: list[Path] = []
    if raw_env_path:
        candidates.append(Path(raw_env_path).expanduser())
    candidates.append(Path(__file__).resolve().parent / 'skills' / 'pptx')

    for candidate in candidates:
        if candidate.is_dir():
            return candidate
    return None


def _run_pptx_skill_script(skill_root: Path, script_relative_path: str, args: list[str]) -> None:
    script_path = skill_root / script_relative_path
    if not script_path.is_file():
        raise FileNotFoundError(f'PPTX skill script not found: {script_path}')

    command = [sys.executable, str(script_path), *args]
    process = subprocess.run(
        command,
        cwd=str(script_path.parent),
        capture_output=True,
        text=True,
        timeout=120,
    )
    if process.returncode == 0:
        return

    stderr = (process.stderr or '').strip()
    stdout = (process.stdout or '').strip()
    details = stderr or stdout or f'exit code {process.returncode}'
    raise RuntimeError(f'PPTX skill command failed ({script_relative_path}): {details}')


def _build_presentation_pptx_bytes_skill(deck: list[dict[str, Any]]) -> bytes:
    skill_root = _resolve_pptx_skill_root()
    if not skill_root:
        raise FileNotFoundError('PPTX skill directory is not available')

    base_bytes = _build_presentation_pptx_bytes_python_pptx(deck)
    with tempfile.TemporaryDirectory(prefix='mts-router-pptx-') as temp_dir:
        temp_path = Path(temp_dir)
        source_pptx = temp_path / 'source.pptx'
        unpacked_dir = temp_path / 'unpacked'
        output_pptx = temp_path / 'output.pptx'
        source_pptx.write_bytes(base_bytes)

        _run_pptx_skill_script(skill_root, 'scripts/office/unpack.py', [str(source_pptx), str(unpacked_dir)])
        _run_pptx_skill_script(skill_root, 'scripts/clean.py', [str(unpacked_dir)])
        _run_pptx_skill_script(
            skill_root,
            'scripts/office/pack.py',
            [str(unpacked_dir), str(output_pptx), '--original', str(source_pptx), '--validate', 'true'],
        )

        if not output_pptx.is_file():
            raise RuntimeError('PPTX skill did not produce output file')
        return output_pptx.read_bytes()


def _render_presentation_pptx(deck: list[dict[str, Any]]) -> tuple[bytes, str]:
    renderer_mode = (os.getenv(PRESENTATION_RENDERER_ENV) or 'auto').strip().lower()
    if renderer_mode in {'python', 'python-pptx', 'legacy'}:
        return _build_presentation_pptx_bytes_python_pptx(deck), 'python-pptx'

    if renderer_mode in {'skill', 'auto'}:
        try:
            return _build_presentation_pptx_bytes_skill(deck), 'pptx-skill'
        except Exception as exc:
            if renderer_mode == 'skill':
                raise
            log.warning(
                'PPTX skill renderer failed, fallback to python-pptx: %s',
                _truncate_error(exc, max_length=360),
            )

    return _build_presentation_pptx_bytes_python_pptx(deck), 'python-pptx'


async def _attach_presentation_file_to_chat(
    request: Request,
    user: UserModel | Any,
    metadata: dict[str, Any],
    deck: list[dict[str, Any]],
    pptx_bytes: bytes,
    renderer: str = 'python-pptx',
    event_emitter: StatusEmitter | None = None,
) -> dict[str, Any] | None:
    if not pptx_bytes:
        return None

    user_id = _get_user_attr(user, 'id')
    if not user_id:
        return None

    first_title = deck[0].get('title') if deck else ''
    file_id = str(uuid.uuid4())
    display_name = _safe_pptx_filename(first_title if isinstance(first_title, str) else '')
    storage_name = f'{file_id}_{display_name}'

    contents, file_path = await asyncio.to_thread(
        Storage.upload_file,
        io.BytesIO(pptx_bytes),
        storage_name,
        {
            'OpenWebUI-User-Email': _get_user_attr(user, 'email'),
            'OpenWebUI-User-Id': user_id,
            'OpenWebUI-User-Name': _get_user_attr(user, 'name'),
            'OpenWebUI-File-Id': file_id,
        },
    )

    file_item = Files.insert_new_file(
        user_id,
        FileForm(
            **{
                'id': file_id,
                'filename': display_name,
                'path': file_path,
                'data': {},
                'meta': {
                    'name': display_name,
                    'content_type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    'size': len(contents),
                    'data': {
                        'generated_by': 'mts-router',
                        'kind': 'presentation',
                        'slides': len(deck),
                        'renderer': renderer,
                    },
                },
            }
        ),
    )
    if not file_item:
        return None

    attachment = {
        'type': 'file',
        'id': file_item.id,
        'url': file_item.id,
        'name': (file_item.meta or {}).get('name') or file_item.filename,
        'size': (file_item.meta or {}).get('size'),
        'content_type': (file_item.meta or {}).get('content_type'),
        'download_url': f'/api/v1/files/{file_item.id}/content',
    }

    chat_id = metadata.get('chat_id') if isinstance(metadata, dict) else None
    message_id = metadata.get('message_id') if isinstance(metadata, dict) else None

    files_payload = [attachment]
    if chat_id and message_id:
        db_files = Chats.add_message_files_by_id_and_message_id(chat_id, message_id, files_payload)
        if db_files:
            files_payload = db_files

    if event_emitter is not None:
        try:
            await event_emitter({'type': 'chat:message:files', 'data': {'files': files_payload}})
        except Exception as emit_exc:
            log.debug(f'Failed to emit presentation file attachment event: {emit_exc}')

    return attachment


def _normalized_fingerprint(text: str) -> str:
    normalized = re.sub(r'\s+', ' ', (text or '').strip())
    return hashlib.sha256(normalized.encode('utf-8')).hexdigest()


def _extract_source_chunks(
    sources: list[dict[str, Any]],
    max_chunks: int = 10,
    max_chars_per_chunk: int = 1200,
    skip_full_markdown: bool = False,
) -> list[dict[str, Any]]:
    chunks: list[dict[str, Any]] = []
    for source in sources or []:
        docs = source.get('document') or []
        metas = source.get('metadata') or []
        distances = source.get('distances') or []
        source_info = source.get('source') or {}
        source_name = source_info.get('name') or source_info.get('id') or 'Unknown'

        for idx, doc in enumerate(docs):
            if len(chunks) >= max_chunks:
                return chunks
            if not isinstance(doc, str) or not doc.strip():
                continue

            metadata = metas[idx] if idx < len(metas) and isinstance(metas[idx], dict) else {}
            if skip_full_markdown and metadata.get('context_mode') == 'full_markdown':
                continue
            distance = distances[idx] if idx < len(distances) else None
            chunk_source = metadata.get('source') or metadata.get('name') or source_name
            chunks.append(
                {
                    'source': str(chunk_source),
                    'distance': distance,
                    'content': doc.strip()[:max_chars_per_chunk],
                }
            )
    return chunks


def _format_chunks_for_prompt(chunks: list[dict[str, Any]], max_total_chars: int = 28_000) -> str:
    if not chunks:
        return ''

    blocks = []
    used_chars = 0
    for index, chunk in enumerate(chunks, start=1):
        header = f"[chunk #{index}] source={chunk.get('source')}"
        if chunk.get('distance') is not None:
            header += f" distance={chunk.get('distance')}"
        body = chunk.get('content') or ''
        block = f'{header}\n{body}'
        if used_chars + len(block) > max_total_chars and blocks:
            break
        blocks.append(block)
        used_chars += len(block)
    return '\n\n'.join(blocks)


def _count_text_tokens(request: Request, text: str) -> int:
    if not text:
        return 0
    encoding_name = str(getattr(request.app.state.config, 'TIKTOKEN_ENCODING_NAME', 'cl100k_base'))
    try:
        encoding = tiktoken.get_encoding(encoding_name)
        return len(encoding.encode(text))
    except Exception:
        return max(1, len(text) // 4)


def _extract_full_markdown_context(
    request: Request,
    sources: list[dict[str, Any]],
    *,
    per_url_token_limit: int = 50_000,
    max_docs: int = 3,
    max_total_tokens: int = 120_000,
) -> list[dict[str, Any]]:
    selected: list[dict[str, Any]] = []
    total_tokens = 0

    for source in sources or []:
        source_info = source.get('source') or {}
        source_type = str(source_info.get('type') or '').strip().lower()
        if source_type != 'web_search':
            continue

        docs = source.get('document') or []
        metas = source.get('metadata') or []
        for idx, doc in enumerate(docs):
            if not isinstance(doc, str) or not doc.strip():
                continue
            metadata = metas[idx] if idx < len(metas) and isinstance(metas[idx], dict) else {}
            if metadata.get('context_mode') != 'full_markdown':
                continue

            token_count = metadata.get('token_count')
            try:
                token_count = int(token_count)
            except Exception:
                token_count = _count_text_tokens(request, doc)

            if token_count > per_url_token_limit:
                continue
            if selected and total_tokens + token_count > max_total_tokens:
                continue

            selected.append(
                {
                    'source': str(metadata.get('title') or metadata.get('source') or source_info.get('name') or 'web'),
                    'url': str(metadata.get('url') or source_info.get('id') or '').strip(),
                    'token_count': token_count,
                    'content': doc.strip(),
                }
            )
            total_tokens += token_count
            if len(selected) >= max_docs:
                return selected

    return selected


def _format_full_markdown_context(
    contexts: list[dict[str, Any]],
    *,
    max_total_chars: int = 300_000,
) -> str:
    if not contexts:
        return ''

    blocks: list[str] = []
    used_chars = 0
    for index, item in enumerate(contexts, start=1):
        header = (
            f"[web full #{index}] source={item.get('source')} "
            f"url={item.get('url')} tokens={item.get('token_count')}"
        )
        body = str(item.get('content') or '')
        remaining = max_total_chars - used_chars
        if remaining <= 0:
            break

        block = f'{header}\n{body}'
        if len(block) <= remaining:
            blocks.append(block)
            used_chars += len(block)
            continue

        # Hard-cap output size even for the first oversized page.
        head = f'{header}\n'
        if len(head) >= remaining:
            break

        allowed_body = max(0, remaining - len(head))
        if allowed_body == 0:
            break

        blocks.append(f'{head}{body[:allowed_body]}')
        used_chars += len(blocks[-1])
        break
    return '\n\n'.join(blocks)


async def _emit_router_status(
    status_emitter: StatusEmitter | None,
    description: str,
    *,
    done: bool = False,
    error: bool = False,
    hidden: bool = False,
    **extra: Any,
) -> None:
    log.info(description)

    if status_emitter is None:
        return

    payload = {
        'type': 'status',
        'data': {
            'action': 'mts_router',
            'description': description,
            'done': done,
            **extra,
        },
    }
    if error:
        payload['data']['error'] = True
    if hidden:
        payload['data']['hidden'] = True

    try:
        await status_emitter(payload)
    except Exception as emit_exc:
        log.debug(f'Failed to emit MTS Router status event: {emit_exc}')


async def _execute_generate_text(
    request: Request,
    form_data: dict[str, Any],
    operation,
    artifact_map: dict[str, Artifact],
    model_id: str,
    previous_text: str | None = None,
    attempt: int = 1,
    max_rounds: int = 1,
) -> Artifact:
    prompt = operation.prompt or _get_text(artifact_map, 'inp_text_1')
    options = dict(getattr(operation, 'options', {}) or {})
    stage = options.get('stage', '')
    agent_role = str(options.get('agent') or '').strip().lower()
    deep_research = _to_bool(options.get('deep_research'))
    is_presentation_stage = str(stage).startswith('presentation_')
    presentation_requirements = (
        _extract_presentation_requirements(artifact_map, prompt)
        if is_presentation_stage
        else {}
    )

    use_passthrough_messages = (
        operation.inputs == ['inp_text_1']
        and form_data.get('messages')
        and attempt == 1
        and not deep_research
    )

    if use_passthrough_messages:
        messages = form_data.get('messages', [])
    else:
        context = _build_context_block(artifact_map, operation.inputs)
        stage_instruction = ''
        if stage == 'draft':
            stage_instruction = 'Сформируй черновик ответа, сохраняя ключевые факты.'
        elif stage == 'final':
            stage_instruction = 'Собери финальный ответ и исправь найденные проблемы.'
        elif stage == 'research_plan':
            stage_instruction = (
                'Собери план исследования: цели, шаги, критерии достаточности фактов и список открытых вопросов.'
            )
        elif stage == 'source_routing':
            stage_instruction = (
                'Раздели задачу на аналитические компоненты и определи, что собирать из интернета, '
                'а что из приложенных документов.'
            )
        elif stage == 'presentation_plan':
            target_slides = int(presentation_requirements.get('slides_target') or 10)
            stage_instruction = (
                'Спроектируй презентацию: цель, целевая аудитория, ключевые тезисы, '
                'ожидаемый формат и требования к качеству контента. '
                f'Целевое число слайдов: {target_slides}.'
            )
        elif stage == 'presentation_outline':
            target_slides = int(presentation_requirements.get('slides_target') or 10)
            required_components = presentation_requirements.get('required_components') or list(PRESENTATION_CORE_COMPONENTS)
            stage_instruction = (
                'Построй детальный outline презентации по слайдам: '
                'роль слайда, главный тезис, факты/цифры, интерпретация и практический вывод. '
                f'Сформируй ровно {target_slides} слайдов и покрой компоненты: {", ".join(required_components)}.'
            )
        elif stage == 'presentation_content':
            target_slides = int(presentation_requirements.get('slides_target') or 10)
            required_components = presentation_requirements.get('required_components') or list(PRESENTATION_CORE_COMPONENTS)
            required_years = presentation_requirements.get('required_years') or []
            period_hint = f'Ключевой период: {", ".join(required_years)}. ' if required_years else ''
            stage_instruction = (
                'Сгенерируй контент слайдов на основе evidence: строго факты, '
                'понятные формулировки, без воды и без выдуманных данных. '
                f'{period_hint}'
                f'Ровно {target_slides} слайдов. '
                f'Обязательные компоненты: {", ".join(required_components)}. '
                'На слайдах с цифрами обязательно указывай период, единицы измерения и источник. '
                'Каждый слайд должен отвечать на 3 вопроса: что произошло, почему это важно, что из этого следует.'
            )
        elif stage == 'presentation_content_refined':
            target_slides = int(presentation_requirements.get('slides_target') or 10)
            required_components = presentation_requirements.get('required_components') or list(PRESENTATION_CORE_COMPONENTS)
            stage_instruction = (
                'Доработай контент слайдов после validation и gap-fill: '
                'исправь ошибки, закрой пробелы и усили аргументацию на слайдах. '
                f'Сохрани ровно {target_slides} слайдов и полное покрытие компонентов: {", ".join(required_components)}.'
            )
        elif stage == 'presentation_finalize':
            target_slides = int(presentation_requirements.get('slides_target') or 10)
            stage_instruction = (
                'Собери финальный пакет презентации: '
                '1) краткий аналитический вывод, '
                '2) рекомендации по использованию deck, '
                '3) финальный JSON со слайдами, '
                '4) ограничения и риски, '
                '5) источники. '
                f'Финальный deck должен содержать ровно {target_slides} слайдов. '
                'Даже при неполном контексте нельзя отказываться от генерации: '
                'верни лучший возможный deck на базе имеющихся источников и явно пометь пробелы.'
            )
        elif stage == 'analysis':
            stage_instruction = (
                'Сделай аналитический синтез по evidence: отдельно по компонентам, '
                'укажи выводы, противоречия, риски и уровень уверенности.'
            )
        elif stage == 'analysis_refined':
            stage_instruction = (
                'Доработай анализ после gap-fill retrieval: закрой пробелы, '
                'обнови цифры и обозначь, какие выводы остались ограниченными.'
            )
        elif stage == 'report':
            stage_instruction = (
                'Собери большой аналитический финальный отчёт с компонентной структурой: '
                '1) Executive Summary (5-8 пунктов), '
                '2) Карта компонентов исследования (что покрыто интернетом и документами), '
                '3) Детальный анализ по каждому компоненту с фактами, цифрами и сравнениями, '
                '4) Сценарии/варианты действий и trade-offs, '
                '5) Ограничения, риски и уровень уверенности по компонентам, '
                '6) Список источников и что именно из них использовано.'
            )

        refinement_instruction = ''
        if previous_text:
            refinement_instruction = (
                'Твоя предыдущая версия уже есть в контексте. Улучши её: убери повторы, '
                'исправь неточности и усили структуру.'
            )

        role_instruction = ''
        if agent_role == 'researcher':
            role_instruction = 'Роль Researcher: оркестрируй исследование и проектируй исполнимый план.'
        elif agent_role == 'analyst':
            role_instruction = 'Роль Analyst: опирайся только на evidence, не добавляй непроверенные факты.'
        elif agent_role == 'reporter':
            role_instruction = 'Роль Reporter: отвечай как автор финального evidence-bound отчёта.'
        elif agent_role == 'reviewer':
            role_instruction = 'Роль Validator: проверяй покрытие и фактическую состоятельность, фиксируй пробелы.'

        research_instruction = ''
        if deep_research:
            research_instruction = (
                'Опирайся на retrieved/chunk-контекст, не выдумывай источники. '
                'Если контекст недостаточен, явно напиши об этом. '
                'Дай подробный ответ, а не краткую выжимку.'
            )
        if is_presentation_stage:
            research_instruction = (
                f'{research_instruction} '
                'Сфокусируйся на доказательной аналитике для слайдов: '
                'приоритет числовым данным, периоду, сравнениям и проверяемым URL-источникам.'
            ).strip()

        structured_output_instruction = ''
        if stage == 'source_routing':
            structured_output_instruction = (
                'Верни ТОЛЬКО JSON-объект без пояснений с полями: '
                '{"needs_internet": boolean, "needs_documents": boolean, '
                '"internet_queries": string[], "document_queries": string[], '
                '"analysis_components": string[]}.'
            )
        elif stage == 'presentation_plan':
            target_slides = int(presentation_requirements.get('slides_target') or 10)
            structured_output_instruction = (
                'Верни ТОЛЬКО JSON-объект без пояснений: '
                '{"objective": string, "audience": string, "tone": string, "language": string, '
                '"slides_target": number, "key_questions": string[], "quality_criteria": string[]}. '
                f'Используй slides_target={target_slides}, если пользователь явно указал число слайдов.'
            )
        elif stage == 'presentation_outline':
            target_slides = int(presentation_requirements.get('slides_target') or 10)
            required_components = presentation_requirements.get('required_components') or list(PRESENTATION_CORE_COMPONENTS)
            structured_output_instruction = (
                'Верни ТОЛЬКО JSON-массив слайдов без пояснений. '
                'Каждый элемент: {"title": string, "component": string, "goal": string, "key_points": string[], '
                '"core_fact": string, "why_it_matters": string, "use_list": boolean}. '
                f'Верни ровно {target_slides} элементов. Покрой компоненты: {", ".join(required_components)}.'
            )
        elif stage in {'presentation_content', 'presentation_content_refined'}:
            target_slides = int(presentation_requirements.get('slides_target') or 10)
            required_components = presentation_requirements.get('required_components') or list(PRESENTATION_CORE_COMPONENTS)
            required_years = presentation_requirements.get('required_years') or []
            period_hint = f'Приоритетный период: {", ".join(required_years)}. ' if required_years else ''
            structured_output_instruction = (
                'Верни ТОЛЬКО JSON-массив слайдов без пояснений. '
                'Каждый элемент: {"title": string, "text": string, "list": string[], '
                '"speaker_notes": string, "image_query": string, "sources": string[], '
                '"key_message": string, "implication": string}. '
                'Поле "list" может быть пустым массивом, если не требуется. '
                f'Ровно {target_slides} слайдов. '
                f'Обязательно покрой компоненты: {", ".join(required_components)}. '
                f'{period_hint}'
                'Требования качества: '
                '1) на содержательных слайдах 2-6 bullet-пунктов; '
                '2) для числовых тезисов добавляй период и единицы измерения; '
                '3) на каждом содержательном слайде хотя бы 1 валидный URL в sources; '
                '4) не пиши общие фразы вроде "данные отсутствуют" без уточнения, что именно не найдено и где искали; '
                '5) поле text обязательно: 3-6 предложений аналитики (факт -> интерпретация -> вывод/следствие), '
                'а не просто перечисление пунктов.'
            )
        elif stage == 'presentation_finalize':
            target_slides = int(presentation_requirements.get('slides_target') or 10)
            structured_output_instruction = (
                'Верни Markdown-ответ со структурой: '
                '1) "## Краткий вывод", '
                '2) "## Рекомендации по подаче", '
                '3) "## Финальный deck JSON" (внутри один json code block), '
                '4) "## Ограничения и риски", '
                '5) "## Источники". '
                f'В "Финальный deck JSON" должен быть только один JSON-массив из {target_slides} слайдов. '
                'Слайды должны быть информативными и опираться на проверяемые источники. '
                'Запрещено отвечать отказом или писать, что презентацию нельзя сформировать.'
            )

        context_payload = [
            f'Запрос пользователя:\n{prompt}'.strip(),
            f'Контекст:\n{context}'.strip() if context else '',
            f'Текущая дата: {time.strftime("%Y-%m-%d")} (используй абсолютные даты, без "сегодня/вчера").',
        ]
        if previous_text:
            context_payload.append(f'Предыдущая версия ответа:\n{previous_text}')
        if stage_instruction:
            context_payload.append(f'Стадия:\n{stage_instruction}')
        if role_instruction:
            context_payload.append(f'Роль:\n{role_instruction}')
        if refinement_instruction:
            context_payload.append(f'Инструкция по доработке:\n{refinement_instruction}')
        if research_instruction:
            context_payload.append(f'Инструкция по верификации:\n{research_instruction}')
        if structured_output_instruction:
            context_payload.append(f'Формат вывода:\n{structured_output_instruction}')
        if max_rounds > 1:
            context_payload.append(f'Раунд генерации: {attempt}/{max_rounds}')

        system_content = (
            'Ты MTS Router. Используй только переданный контекст и ответь по делу. '
            'Если данных недостаточно, честно скажи об этом.'
        )
        if is_presentation_stage:
            system_content = (
                f'{system_content} '
                'Для презентационных этапов всегда формируй результат по доступным данным и не отказывайся '
                'из-за отсутствия прямого доступа к интернету: контекст уже предоставлен выше.'
            )

        messages = [
            {
                'role': 'system',
                'content': system_content,
            },
            {
                'role': 'user',
                'content': '\n\n'.join([block for block in context_payload if block]).strip(),
            },
        ]

    temperature = 0.2
    if stage in {'source_routing', 'research_plan', 'presentation_plan', 'presentation_outline'}:
        temperature = 0
    elif deep_research or stage.startswith('presentation_'):
        temperature = 0.1

    text = await create_text_response(
        request=request,
        model=model_id,
        messages=messages,
        temperature=temperature,
    )
    metadata = {'model_id': model_id}
    if stage:
        metadata['stage'] = stage
    if agent_role:
        metadata['agent_role'] = agent_role
    if max_rounds > 1:
        metadata['attempt'] = attempt
        metadata['max_rounds'] = max_rounds
    return Artifact(id=operation.output_id, type='text', source='operation', text=text, metadata=metadata)


async def _execute_analyze_image(request: Request, operation, artifact_map: dict[str, Artifact], model_id: str) -> Artifact:
    image_artifact = next((artifact_map.get(input_id) for input_id in operation.inputs if artifact_map.get(input_id)), None)
    if not image_artifact or not image_artifact.data:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail='Image input is missing for the analyze_image operation.',
        )

    messages = [
        {
            'role': 'system',
            'content': (
                'Ты анализируешь изображение для последующей оркестрации. '
                'Опиши важные визуальные факты, видимый текст и сущности компактно и без домыслов.'
            ),
        },
        {
            'role': 'user',
            'content': [
                {'type': 'text', 'text': operation.prompt or 'Проанализируй изображение.'},
                {'type': 'image_url', 'image_url': {'url': image_artifact.data}},
            ],
        },
    ]
    text = await create_text_response(
        request=request,
        model=model_id,
        messages=messages,
        temperature=0,
    )
    return Artifact(
        id=operation.output_id,
        type='structured',
        source='operation',
        text=text,
        metadata={'model_id': model_id, 'source_artifact_id': image_artifact.id},
    )


async def _execute_transcribe_audio(request: Request, operation, artifact_map: dict[str, Artifact], model_id: str) -> Artifact:
    audio_artifact = next((artifact_map.get(input_id) for input_id in operation.inputs if artifact_map.get(input_id)), None)
    if not audio_artifact or not audio_artifact.path:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail='Audio input is missing for the transcribe_audio operation.',
        )

    text = await transcribe_audio(
        request=request,
        model=model_id,
        file_path=audio_artifact.path,
        language=(audio_artifact.metadata or {}).get('language'),
    )
    return Artifact(
        id=operation.output_id,
        type='text',
        source='operation',
        text=text,
        metadata={'model_id': model_id, 'source_artifact_id': audio_artifact.id},
    )


def _to_positive_int(value: Any, default: int) -> int:
    try:
        parsed = int(value)
        if parsed > 0:
            return parsed
    except Exception:
        pass
    return default


def _to_bool(value: Any) -> bool:
    if isinstance(value, str):
        return value.strip().lower() in {'1', 'true', 'yes', 'on'}
    return bool(value)


def _normalize_query_list(value: Any, max_items: int = 6) -> list[str]:
    if value is None:
        return []

    if isinstance(value, str):
        raw_items = [value]
    elif isinstance(value, list):
        raw_items = value
    else:
        return []

    normalized: list[str] = []
    seen: set[str] = set()
    for item in raw_items:
        candidate = re.sub(r'\s+', ' ', str(item or '').strip())
        if not candidate:
            continue
        lowered = candidate.lower()
        if lowered in seen:
            continue
        seen.add(lowered)
        normalized.append(candidate[:240])
        if len(normalized) >= max_items:
            break
    return normalized


def _merge_query_lists(*groups: Any, max_items: int = 6) -> list[str]:
    merged: list[str] = []
    seen: set[str] = set()
    for group in groups:
        for item in _normalize_query_list(group, max_items=max_items):
            lowered = item.lower()
            if lowered in seen:
                continue
            seen.add(lowered)
            merged.append(item)
            if len(merged) >= max_items:
                return merged
    return merged


def _extract_json_payload(text: str) -> dict[str, Any]:
    if not text:
        return {}

    payload = text.strip()
    candidates: list[str] = []
    if payload:
        candidates.append(payload)

    fenced = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', text, flags=re.IGNORECASE | re.DOTALL)
    if fenced:
        candidates.insert(0, fenced.group(1).strip())

    start_index = text.find('{')
    end_index = text.rfind('}')
    if start_index != -1 and end_index != -1 and end_index > start_index:
        candidates.append(text[start_index : end_index + 1].strip())

    for candidate in candidates:
        try:
            parsed = json.loads(candidate)
            if isinstance(parsed, dict):
                return parsed
        except Exception:
            continue
    return {}


def _extract_requested_slide_count(text: str, default: int = 10) -> int:
    if not text:
        return default
    candidates = re.findall(
        r'(\d{1,2})\s*(?:слайд(?:а|ов)?|slides?)',
        text,
        flags=re.IGNORECASE,
    )
    if not candidates:
        return default
    try:
        value = int(candidates[0])
    except Exception:
        return default
    return max(4, min(20, value))


def _extract_requested_years(text: str) -> list[str]:
    if not text:
        return []
    years = re.findall(r'\b(20\d{2})\b', text)
    ordered: list[str] = []
    seen: set[str] = set()
    for year in years:
        if year in seen:
            continue
        seen.add(year)
        ordered.append(year)
    return ordered[:3]


def _extract_presentation_subject(query: str) -> str:
    cleaned = re.sub(r'\s+', ' ', (query or '').strip())
    if not cleaned:
        return ''

    # Common command pattern: "Сделай презентацию ...: <subject>. Нужны: ..."
    if ':' in cleaned:
        _, _, tail = cleaned.partition(':')
        if tail.strip():
            cleaned = tail.strip()

    cleaned = re.split(r'\bнужн[аоы]\s*:?', cleaned, flags=re.IGNORECASE)[0].strip(' .;,-')
    cleaned = re.sub(r'^(сделай|подготовь|собери)\s+презентац[июя]\b', '', cleaned, flags=re.IGNORECASE).strip(' .;,-')
    return cleaned or query.strip()


def _extract_presentation_requirements(
    artifact_map: dict[str, Artifact],
    prompt: str,
) -> dict[str, Any]:
    slides_target = _extract_requested_slide_count(prompt, default=10)
    years = _extract_requested_years(prompt)
    components: list[str] = list(PRESENTATION_CORE_COMPONENTS)
    quality_criteria: list[str] = []

    plan_payload: dict[str, Any] = {}
    for artifact in artifact_map.values():
        stage = str((artifact.metadata or {}).get('stage') or '').strip().lower()
        if stage != 'presentation_plan' or not artifact.text:
            continue
        parsed = _extract_json_payload(artifact.text)
        if parsed:
            plan_payload = parsed
            break

    plan_slides = plan_payload.get('slides_target')
    try:
        if plan_slides is not None:
            slides_target = max(4, min(20, int(plan_slides)))
    except Exception:
        pass

    quality_criteria = _normalize_query_list(plan_payload.get('quality_criteria'), max_items=8)
    key_questions = _normalize_query_list(plan_payload.get('key_questions'), max_items=8)
    plan_text = '\n'.join([str(plan_payload.get('objective') or ''), *key_questions, *quality_criteria])
    years = _merge_query_lists(years, _extract_requested_years(plan_text), max_items=3)

    lowered_prompt = (prompt or '').lower()
    optional_components = (
        ('динамика', ('динамик', 'рост', 'снижен', 'изменен')),
        ('ставки', ('ставк', 'доходност')),
        ('конкуренция банков', ('банки', 'игрок', 'топ-')),
        ('сценарии', ('сценар', 'прогноз')),
        ('рекомендации', ('рекомендац', 'что делать')),
    )
    for component_name, markers in optional_components:
        if any(marker in lowered_prompt for marker in markers):
            components.append(component_name)

    components = _merge_query_lists(components, max_items=10)
    return {
        'slides_target': slides_target,
        'required_components': components or list(PRESENTATION_CORE_COMPONENTS),
        'required_years': years,
        'quality_criteria': quality_criteria,
        'subject': _extract_presentation_subject(prompt),
    }


def _augment_presentation_queries(
    query: str,
    document_queries: list[str],
    web_queries: list[str],
    *,
    max_items: int = 8,
) -> tuple[list[str], list[str]]:
    subject = _extract_presentation_subject(query)
    years = _extract_requested_years(query)
    year = years[0] if years else ''
    period_suffix = f' {year}' if year else ''
    subject_with_period = f'{subject}{period_suffix}'.strip()

    web_templates = [
        subject_with_period,
        f'ключевые цифры {subject_with_period}'.strip(),
        f'структура рынка {subject_with_period}'.strip(),
        f'риски {subject_with_period}'.strip(),
        f'Банк России {subject_with_period} статистика'.strip(),
        f'официальная статистика {subject_with_period} site:cbr.ru'.strip(),
        f'объем вкладов населения {period_suffix} site:cbr.ru'.strip(),
    ]

    document_templates = [
        subject_with_period,
        f'структура {subject_with_period}'.strip(),
        f'ключевые показатели {subject_with_period}'.strip(),
        f'риски {subject_with_period}'.strip(),
        f'выводы {subject_with_period}'.strip(),
    ]

    merged_docs = _merge_query_lists(document_queries, document_templates, max_items=max_items)
    merged_web = _merge_query_lists(web_queries, web_templates, max_items=max_items)
    return merged_docs, merged_web


def _extract_source_routing_plan(artifact_map: dict[str, Artifact], input_ids: list[str], base_query: str) -> dict[str, Any]:
    defaults = {
        'needs_internet': True,
        'needs_documents': True,
        'internet_queries': [base_query] if base_query else [],
        'document_queries': [base_query] if base_query else [],
        'analysis_components': [],
    }

    routing_artifact = None
    for input_id in input_ids or []:
        artifact = artifact_map.get(input_id)
        if not artifact:
            continue
        stage = str((artifact.metadata or {}).get('stage') or '').strip().lower()
        if stage == 'source_routing':
            routing_artifact = artifact
            break

    if not routing_artifact or not routing_artifact.text:
        return defaults

    parsed = _extract_json_payload(routing_artifact.text)
    if not parsed:
        return defaults

    needs_internet = _to_bool(parsed.get('needs_internet', defaults['needs_internet']))
    needs_documents = _to_bool(parsed.get('needs_documents', defaults['needs_documents']))
    internet_queries = _merge_query_lists(parsed.get('internet_queries'), base_query, max_items=6)
    document_queries = _merge_query_lists(parsed.get('document_queries'), base_query, max_items=6)
    components = _normalize_query_list(parsed.get('analysis_components'), max_items=10)

    return {
        'needs_internet': needs_internet,
        'needs_documents': needs_documents,
        'internet_queries': internet_queries,
        'document_queries': document_queries,
        'analysis_components': components,
    }


def _extract_validation_queries(artifact_map: dict[str, Artifact], input_ids: list[str], base_query: str) -> list[str]:
    validation_text = ''
    for input_id in input_ids or []:
        artifact = artifact_map.get(input_id)
        if not artifact:
            continue
        stage = str((artifact.metadata or {}).get('stage') or '').strip().lower()
        if stage in {'validation', 'review', 'presentation_validation'} and artifact.text:
            validation_text = artifact.text
            break

    if not validation_text:
        return []

    parsed = _extract_json_payload(validation_text)
    queries = _merge_query_lists(
        parsed.get('internet_queries'),
        parsed.get('followup_queries'),
        parsed.get('document_queries'),
        max_items=6,
    )
    if queries:
        return queries

    hints: list[str] = []
    for raw_line in validation_text.splitlines():
        line = re.sub(r'^[\-\*\d\.\)\s]+', '', raw_line or '').strip()
        if not line:
            continue
        lowered = line.lower()
        if any(
            token in lowered
            for token in (
                'недостат',
                'пробел',
                'missing',
                'не хватает',
                'не покрыт',
                'уточнить',
                'нужно',
                'нет данных',
            )
        ):
            hints.append(line[:180])
        if len(hints) >= 4:
            break

    return _merge_query_lists([base_query], hints, max_items=6)


def _build_deep_research_queries(query: str, max_queries: int = 3) -> list[str]:
    cleaned = re.sub(r'\s+', ' ', (query or '').strip())
    if not cleaned:
        return []

    candidates = [cleaned]
    lowered = cleaned.lower()
    words = cleaned.split()

    if len(words) > 7:
        candidates.append(' '.join(words[:8]))

    if any(
        token in lowered
        for token in (
            'api',
            'sdk',
            'документ',
            'документац',
            'model',
            'модель',
            'error',
            'ошиб',
            'docker',
            'python',
        )
    ):
        candidates.append(f'{cleaned} official documentation')

    unique: list[str] = []
    seen: set[str] = set()
    for item in candidates:
        normalized = item.strip()
        lowered_item = normalized.lower()
        if not normalized or lowered_item in seen:
            continue
        seen.add(lowered_item)
        unique.append(normalized)
        if len(unique) >= max_queries:
            break
    return unique


async def _collect_web_research_items(
    request: Request,
    user: UserModel | Any,
    query: str,
    max_results: int,
    seed_queries: list[str] | None = None,
    max_queries: int = 6,
    status_emitter: StatusEmitter | None = None,
    status_context: dict[str, Any] | None = None,
) -> list[dict[str, Any]]:
    if not bool(getattr(request.app.state.config, 'ENABLE_WEB_SEARCH', False)):
        return []

    normalized_max_queries = max(2, min(10, int(max_queries or 6)))
    seed = _merge_query_lists(seed_queries or [], query, max_items=normalized_max_queries)
    expanded: list[str] = []
    for item in seed:
        expanded.extend(_build_deep_research_queries(item, max_queries=2))
    queries = _merge_query_lists(seed, expanded, max_items=normalized_max_queries)
    if not queries:
        return []

    context = dict(status_context or {})
    await _emit_router_status(
        status_emitter,
        'MTS Router: Searcher — сформировал web-поисковые запросы',
        done=False,
        phase='research',
        action='web_search_queries_generated',
        queries=queries,
        **context,
    )

    configured_engine = str(getattr(request.app.state.config, 'WEB_SEARCH_ENGINE', '') or '').strip()
    engines: list[str] = [configured_engine] if configured_engine else []
    if 'duckduckgo' not in {engine.lower() for engine in engines}:
        engines.append('duckduckgo')

    items: list[dict[str, Any]] = []
    seen_urls: set[str] = set()

    for query_item in queries:
        query_results = None
        selected_engine = None

        for engine in engines:
            try:
                results = await asyncio.to_thread(retrieval_search_web, request, engine, query_item, user)
                if results:
                    query_results = results
                    selected_engine = engine
                    break
            except Exception as exc:
                log.debug(f'Web search failed for engine={engine}, query={query_item}: {exc}')

        if not query_results:
            continue

        query_items_preview: list[dict[str, str]] = []
        for rank, result in enumerate(query_results, start=1):
            url = str(getattr(result, 'link', '') or '').strip()
            if not url or url in seen_urls:
                continue
            seen_urls.add(url)
            title = str(getattr(result, 'title', '') or '').strip()
            items.append(
                {
                    'query': query_item,
                    'engine': selected_engine or configured_engine or 'unknown',
                    'rank': rank,
                    'url': url,
                    'title': title,
                    'snippet': str(getattr(result, 'snippet', '') or '').strip(),
                }
            )
            if len(query_items_preview) < 8:
                query_items_preview.append({'link': url, 'title': title or url})
            if len(items) >= max_results:
                await _emit_router_status(
                    status_emitter,
                    f'MTS Router: Searcher — web-поиск по "{query_item}" дал {len(query_items_preview)} результатов',
                    done=False,
                    phase='research',
                    action='web_search',
                    query=query_item,
                    items=query_items_preview,
                    search_engine=selected_engine or configured_engine or 'unknown',
                    **context,
                )
                return items

        if query_items_preview:
            await _emit_router_status(
                status_emitter,
                f'MTS Router: Searcher — web-поиск по "{query_item}" дал {len(query_items_preview)} результатов',
                done=False,
                phase='research',
                action='web_search',
                query=query_item,
                items=query_items_preview,
                search_engine=selected_engine or configured_engine or 'unknown',
                **context,
            )

    return items


async def _hydrate_web_items(
    request: Request,
    items: list[dict[str, Any]],
    max_pages: int,
    max_chars_per_page: int = 1800,
) -> list[dict[str, Any]]:
    if not items:
        return []

    semaphore = asyncio.Semaphore(3)

    async def _load(item: dict[str, Any]) -> dict[str, Any] | None:
        async with semaphore:
            url = item.get('url', '')
            page_text = ''

            try:
                content, _ = await asyncio.to_thread(get_content_from_url, request, url)
                page_text = re.sub(r'\n{3,}', '\n\n', str(content or '')).strip()
            except Exception as exc:
                log.debug(f'Failed to fetch web page content for {url}: {exc}')

            if not page_text:
                page_text = (item.get('snippet') or '').strip()

            if not page_text:
                return None

            token_count = _count_text_tokens(request, page_text)
            full_markdown = token_count <= 50_000
            prepared_content = page_text if full_markdown else page_text[:max_chars_per_page]

            return {
                **item,
                'content': prepared_content,
                'token_count': token_count,
                'context_mode': 'full_markdown' if full_markdown else 'chunked',
            }

    tasks = [_load(item) for item in items[:max_pages]]
    resolved = await asyncio.gather(*tasks)
    return [item for item in resolved if item]


def _web_items_to_sources(items: list[dict[str, Any]]) -> list[dict[str, Any]]:
    sources: list[dict[str, Any]] = []
    for item in items:
        url = str(item.get('url') or '').strip()
        content = str(item.get('content') or '').strip()
        if not url or not content:
            continue

        title = str(item.get('title') or '').strip()
        source_name = title or urlparse(url).netloc or url
        sources.append(
            {
                'source': {
                    'name': source_name,
                    'id': url,
                    'type': 'web_search',
                },
                'document': [content],
                'metadata': [
                    {
                        'source': url,
                        'url': url,
                        'title': title,
                        'query': item.get('query'),
                        'engine': item.get('engine'),
                        'rank': item.get('rank'),
                        'token_count': item.get('token_count'),
                        'context_mode': item.get('context_mode'),
                        'kind': 'web_search',
                    }
                ],
                'distances': [None],
            }
        )
    return sources


def _dedupe_sources(sources: list[dict[str, Any]]) -> list[dict[str, Any]]:
    deduped: list[dict[str, Any]] = []
    seen: set[str] = set()
    for source in sources:
        source_info = source.get('source') or {}
        source_id = str(source_info.get('id') or source_info.get('name') or '')
        first_doc = ''
        documents = source.get('document') or []
        if documents and isinstance(documents[0], str):
            first_doc = documents[0][:500]

        fingerprint = hashlib.sha1(f'{source_id}::{first_doc}'.encode('utf-8')).hexdigest()
        if fingerprint in seen:
            continue
        seen.add(fingerprint)
        deduped.append(source)
    return deduped


async def _collect_research_sources(
    request: Request,
    form_data: dict[str, Any],
    user: UserModel | Any,
    query: str,
    top_k: int,
    deep_research: bool = False,
    stage: str = '',
    source_routing: dict[str, Any] | None = None,
    followup_queries: list[str] | None = None,
    status_emitter: StatusEmitter | None = None,
    status_context: dict[str, Any] | None = None,
) -> list[dict[str, Any]]:
    metadata = dict(form_data.get('metadata', {}) or {})
    merged_sources: list[dict[str, Any]] = list(metadata.get('sources') or [])
    context = dict(status_context or {})
    routing = dict(source_routing or {})

    include_documents = _to_bool(routing.get('needs_documents', True))
    include_internet = deep_research and _to_bool(routing.get('needs_internet', True))
    document_queries = _merge_query_lists(routing.get('document_queries'), query, max_items=6)
    web_queries = _merge_query_lists(routing.get('internet_queries'), query, max_items=6)
    is_presentation_stage = stage.startswith('presentation_')
    if not document_queries:
        document_queries = _merge_query_lists(query, max_items=1)
    if not web_queries:
        web_queries = _merge_query_lists(query, max_items=1)

    if followup_queries:
        followup = _merge_query_lists(followup_queries, max_items=6)
        document_queries = _merge_query_lists(document_queries, followup, max_items=6)
        web_queries = _merge_query_lists(web_queries, followup, max_items=6)

    if is_presentation_stage:
        document_queries, web_queries = _augment_presentation_queries(
            query=query,
            document_queries=document_queries,
            web_queries=web_queries,
            max_items=8,
        )

    await _emit_router_status(
        status_emitter,
        (
            'MTS Router: Searcher — source-routing '
            f'(stage={stage or "default"}, docs={"on" if include_documents else "off"}, '
            f'web={"on" if include_internet else "off"})'
        ),
        done=False,
        phase='research',
        action='source_routing',
        document_queries=document_queries[:4],
        web_queries=web_queries[:4],
        components=_normalize_query_list(routing.get('analysis_components'), max_items=8),
        **context,
    )

    files = metadata.get('files') or form_data.get('files') or []
    if include_documents and files and request.app.state.EMBEDDING_FUNCTION:
        await _emit_router_status(
            status_emitter,
            f'MTS Router: Searcher — ищу релевантные фрагменты в документах ({len(files)} файлов)',
            done=False,
            phase='research',
            **context,
        )
        try:
            file_sources = await get_sources_from_items(
                request=request,
                items=files,
                queries=document_queries[:4] or [query],
                embedding_function=lambda q, prefix: request.app.state.EMBEDDING_FUNCTION(q, prefix=prefix, user=user),
                k=top_k,
                reranking_function=(
                    (lambda q, docs: request.app.state.RERANKING_FUNCTION(q, docs, user=user))
                    if request.app.state.RERANKING_FUNCTION
                    else None
                ),
                k_reranker=max(10, int(getattr(request.app.state.config, 'TOP_K_RERANKER', 10))),
                r=float(getattr(request.app.state.config, 'RELEVANCE_THRESHOLD', 0.0)),
                hybrid_bm25_weight=float(getattr(request.app.state.config, 'HYBRID_BM25_WEIGHT', 0.5)),
                hybrid_search=bool(getattr(request.app.state.config, 'ENABLE_RAG_HYBRID_SEARCH', False)),
                full_context=False,
                user=user,
            )
            merged_sources.extend(file_sources or [])
            await _emit_router_status(
                status_emitter,
                f'MTS Router: Searcher — из документов извлечено источников: {len(file_sources or [])}',
                done=False,
                phase='research',
                action='sources_retrieved',
                count=len(file_sources or []),
                source_kind='documents',
                **context,
            )
        except Exception as exc:
            log.debug(f'Failed to retrieve research sources via retrieval pipeline: {exc}')
            await _emit_router_status(
                status_emitter,
                f'MTS Router: Searcher — ошибка retrieval по документам: {_truncate_error(exc)}',
                done=False,
                error=True,
                phase='research',
                **context,
            )

    if include_internet:
        await _emit_router_status(
            status_emitter,
            'MTS Router: Searcher — запускаю web research',
            done=False,
            phase='research',
            **context,
        )
        max_web_results = min(max(top_k * 2, 8), 30 if is_presentation_stage else 20)
        max_web_queries = 8 if is_presentation_stage else 6
        web_items = await _collect_web_research_items(
            request=request,
            user=user,
            query=query,
            max_results=max_web_results,
            seed_queries=web_queries[:max_web_queries],
            max_queries=max_web_queries,
            status_emitter=status_emitter,
            status_context=context,
        )
        hydrated_items = await _hydrate_web_items(
            request=request,
            items=web_items,
            max_pages=min(max(top_k, 6), 12 if is_presentation_stage else 10),
        )
        web_sources = _web_items_to_sources(hydrated_items)
        merged_sources.extend(web_sources)
        await _emit_router_status(
            status_emitter,
            f'MTS Router: Searcher — web research добавил источников: {len(web_sources)}',
            done=False,
            phase='research',
            action='sources_retrieved',
            count=len(web_sources),
            source_kind='web',
            **context,
        )

    deduped_sources = _dedupe_sources(merged_sources)
    await _emit_router_status(
        status_emitter,
        f'MTS Router: Searcher — итоговый evidence-пул: {len(deduped_sources)} источников',
        done=False,
        phase='research',
        action='sources_retrieved',
        count=len(deduped_sources),
        source_kind='merged',
        **context,
    )
    return deduped_sources


async def _execute_retrieve_context(
    request: Request,
    form_data: dict[str, Any],
    user: UserModel | Any,
    operation,
    artifact_map: dict[str, Artifact],
    model_id: str,
    status_emitter: StatusEmitter | None = None,
    status_context: dict[str, Any] | None = None,
) -> Artifact:
    prompt = operation.prompt or 'Собери релевантный контекст по запросу.'
    options = dict(getattr(operation, 'options', {}) or {})
    stage = str(options.get('stage') or '').strip().lower()
    agent_role = str(options.get('agent') or '').strip().lower()
    deep_research = _to_bool(options.get('deep_research'))
    top_k_chunks = _to_positive_int(options.get('top_k_chunks'), 20)
    top_k_chunks = min(top_k_chunks, 20)
    context = dict(status_context or {})
    source_routing = _extract_source_routing_plan(artifact_map, operation.inputs, prompt)
    followup_queries = (
        _extract_validation_queries(artifact_map, operation.inputs, prompt)
        if stage in {'gap_fill', 'search_gap', 'presentation_gap_fill'}
        else []
    )
    cache_state = form_data.setdefault('_mts_router_cache', {})
    retrieve_cache = cache_state.setdefault('retrieve_context', {})
    cache_payload = {
        'stage': stage,
        'prompt': prompt,
        'top_k_chunks': top_k_chunks,
        'deep_research': deep_research,
        'source_routing': source_routing,
        'followup_queries': followup_queries,
    }
    cache_key = hashlib.sha1(
        json.dumps(cache_payload, ensure_ascii=False, sort_keys=True, default=str).encode('utf-8')
    ).hexdigest()

    await _emit_router_status(
        status_emitter,
        (
            f'MTS Router: Searcher — собираю контекст '
            f'(top_k={top_k_chunks}, deep_research={"on" if deep_research else "off"})'
        ),
        done=False,
        phase='research',
        **context,
    )
    if followup_queries:
        await _emit_router_status(
            status_emitter,
            'MTS Router: Searcher — validator подсказал уточняющие запросы для gap-fill',
            done=False,
            phase='research',
            action='gap_fill_queries',
            queries=followup_queries[:6],
            **context,
        )

    if cache_key in retrieve_cache:
        sources = retrieve_cache.get(cache_key, [])
        await _emit_router_status(
            status_emitter,
            'MTS Router: Searcher — использую кэшированный evidence-пул (без повторного web-search)',
            done=False,
            phase='research',
            action='cache_hit',
            cache_key=cache_key[:12],
            source_count=len(sources),
            **context,
        )
    else:
        sources = await _collect_research_sources(
            request,
            form_data,
            user,
            prompt,
            top_k_chunks,
            deep_research=deep_research,
            stage=stage,
            source_routing=source_routing,
            followup_queries=followup_queries,
            status_emitter=status_emitter,
            status_context=context,
        )
        retrieve_cache[cache_key] = sources
    context_window_tokens = _context_window_hint_for_model(model_id)
    normalizer_budget_tokens = max(8_000, int(context_window_tokens * 0.5))
    full_markdown_total_tokens_limit = min(50_000, max(6_000, int(normalizer_budget_tokens * 0.7)))
    full_markdown_per_url_tokens_limit = min(20_000, max(3_000, int(full_markdown_total_tokens_limit * 0.6)))
    context_char_budget = max(40_000, normalizer_budget_tokens * 4)
    chunks_char_budget = min(20_000, max(8_000, int(context_char_budget * 0.3)))
    full_context_char_budget = min(120_000, max(30_000, int(context_char_budget * 0.8)))
    full_context_medium_budget = min(full_context_char_budget, max(20_000, int(full_context_char_budget * 0.7)))
    full_context_small_budget = min(full_context_char_budget, max(12_000, int(full_context_char_budget * 0.45)))

    full_markdown_contexts = _extract_full_markdown_context(
        request=request,
        sources=sources,
        per_url_token_limit=full_markdown_per_url_tokens_limit,
        max_total_tokens=full_markdown_total_tokens_limit,
    )
    full_markdown_block = _format_full_markdown_context(full_markdown_contexts, max_total_chars=full_context_char_budget)
    if full_markdown_contexts:
        await _emit_router_status(
            status_emitter,
            (
                f'MTS Router: Searcher — в полный контекст добавлены full-markdown страницы: '
                f'{len(full_markdown_contexts)} '
                f'(лимит: {full_markdown_per_url_tokens_limit} токенов/URL, '
                f'общий: {full_markdown_total_tokens_limit})'
            ),
            done=False,
            phase='research',
            **context,
        )

    chunks = _extract_source_chunks(sources, max_chunks=top_k_chunks, skip_full_markdown=True)
    chunks_block = _format_chunks_for_prompt(chunks, max_total_chars=chunks_char_budget)
    full_docs_count = len(full_markdown_contexts)
    chunk_count = len(chunks)
    context_block = '\n\n'.join([block for block in [full_markdown_block, chunks_block] if block]).strip()

    if not context_block:
        await _emit_router_status(
            status_emitter,
            'MTS Router: Searcher — контекстные источники не найдены',
            done=False,
            phase='research',
            action='sources_retrieved',
            count=0,
            **context,
        )
        return Artifact(
            id=operation.output_id,
            type='structured',
            source='operation',
            text='Контекстные источники не найдены. Продолжай ответ с явным указанием недостатка данных.',
            metadata={'model_id': model_id, 'chunks_used': 0, 'source_count': 0},
        )

    await _emit_router_status(
        status_emitter,
        (
            f'MTS Router: Searcher — нормализую контекст '
            f'(full-markdown={full_docs_count}, chunks={chunk_count})'
        ),
        done=False,
        phase='research',
        **context,
    )

    if agent_role == 'searcher' or stage == 'search':
        system_prompt = (
            'Ты Searcher для MTS Router. Нормализуй evidence без домыслов. '
            'Верни: 1) факты 2) карта источников chunk->source 3) пробелы/конфликты.'
        )
        user_prompt = (
            f'Запрос:\n{prompt}\n\n'
            f'Извлечённый контекст (full-markdown + чанки, используй до {top_k_chunks} чанков):\n{context_block}\n\n'
            'Собери evidence-pack для Analyst/Reporter.'
        )
    else:
        system_prompt = (
            'Ты research-retriever для MTS Router. Сжато нормализуй контекст. '
            'Вывод: 1) ключевые факты 2) список источников по chunk-id.'
        )
        user_prompt = ''

    chunks_compact = _format_chunks_for_prompt(chunks[: min(len(chunks), 8)], max_total_chars=min(chunks_char_budget, 12_000))
    full_markdown_medium = _format_full_markdown_context(
        full_markdown_contexts,
        max_total_chars=full_context_medium_budget,
    )
    full_markdown_small = _format_full_markdown_context(
        full_markdown_contexts,
        max_total_chars=full_context_small_budget,
    )
    context_candidates: list[tuple[str, str]] = []
    seen_candidates: set[str] = set()
    for label, block in [
        ('chunks_only', chunks_block),
        ('full_small_or_chunks_compact', '\n\n'.join([b for b in [full_markdown_small, chunks_compact] if b]).strip()),
        ('full_medium+chunks_compact', '\n\n'.join([b for b in [full_markdown_medium, chunks_compact] if b]).strip()),
        ('full+chunks', context_block),
    ]:
        candidate = (block or '').strip()
        if not candidate:
            continue
        fingerprint = _normalized_fingerprint(candidate)
        if fingerprint in seen_candidates:
            continue
        seen_candidates.add(fingerprint)
        context_candidates.append((label, candidate))

    last_context_error: Exception | None = None
    normalized_context = None
    selected_context_variant = None
    for variant_index, (variant_name, context_candidate) in enumerate(context_candidates, start=1):
        if agent_role == 'searcher' or stage == 'search':
            user_prompt = (
                f'Запрос:\n{prompt}\n\n'
                f'Извлечённый контекст (вариант: {variant_name}, full-markdown + чанки):\n{context_candidate}\n\n'
                'Собери evidence-pack для Analyst/Reporter.'
            )
        else:
            user_prompt = (
                f'Запрос:\n{prompt}\n\n'
                f'Извлечённый контекст (вариант: {variant_name}, full-markdown + чанки):\n{context_candidate}\n\n'
                'Верни краткий структурированный контекст для следующего этапа генерации.'
            )

        messages = [
            {
                'role': 'system',
                'content': system_prompt,
            },
            {
                'role': 'user',
                'content': user_prompt,
            },
        ]
        try:
            normalized_context = await create_text_response(
                request=request,
                model=model_id,
                messages=messages,
                temperature=0,
            )
            selected_context_variant = variant_name
            break
        except Exception as exc:
            last_context_error = exc
            log.warning(
                'MTS Router context normalization failed: variant=%s model=%s err=%s',
                variant_name,
                model_id,
                _truncate_error(exc, max_length=400),
            )
            if variant_index < len(context_candidates):
                await _emit_router_status(
                    status_emitter,
                    'MTS Router: Searcher — уплотняю контекст и повторяю нормализацию',
                    done=False,
                    hidden=True,
                    phase='research',
                    action='context_retry',
                    context_variant=variant_name,
                    **context,
                )
                continue
            raise

    if normalized_context is None:
        raise last_context_error or RuntimeError('Failed to normalize research context.')

    return Artifact(
        id=operation.output_id,
        type='structured',
        source='operation',
        text=normalized_context,
        metadata={
            'model_id': model_id,
            'chunks_used': len(chunks),
            'source_count': len(sources),
            'full_markdown_docs': len(full_markdown_contexts),
            'context_variant': selected_context_variant,
            'stage': stage or None,
            'agent_role': agent_role or None,
        },
    )


async def _execute_critique_text(
    request: Request,
    operation,
    artifact_map: dict[str, Artifact],
    model_id: str,
) -> Artifact:
    prompt = operation.prompt or 'Проверь качество ответа.'
    options = dict(getattr(operation, 'options', {}) or {})
    stage = str(options.get('stage') or '').strip().lower()
    agent_role = str(options.get('agent') or '').strip().lower()
    context = _build_context_block(artifact_map, operation.inputs)

    if stage in {'validation', 'presentation_validation'}:
        if stage == 'presentation_validation':
            validation_schema = (
                '{\n'
                '  "coverage_score": 0-100,\n'
                '  "factual_consistency_score": 0-100,\n'
                '  "structure_score": 0-100,\n'
                '  "numeric_evidence_score": 0-100,\n'
                '  "source_quality_score": 0-100,\n'
                '  "confidence_score": 0-100,\n'
                '  "missing_topics": ["..."],\n'
                '  "contradictions": ["..."],\n'
                '  "slide_level_issues": [\n'
                '    {"slide_title": "...", "issues": ["..."], "fix": "..."}\n'
                '  ],\n'
                '  "missing_numeric_claims": ["..."],\n'
                '  "weak_or_missing_sources": ["..."],\n'
                '  "internet_queries": ["..."],\n'
                '  "document_queries": ["..."],\n'
                '  "followup_queries": ["..."],\n'
                '  "fix_steps": ["..."]\n'
                '}'
            )
            validation_hint = (
                'Сфокусируйся на проверке презентации: корректность фактов, связность повествования между слайдами, '
                'отсутствие повторов и реалистичность рекомендаций. '
                'Отдельно проверь наличие числовых фактов с единицами измерения и качество URL-источников по слайдам.'
            )
        else:
            validation_schema = (
                '{\n'
                '  "coverage_score": 0-100,\n'
                '  "factual_consistency_score": 0-100,\n'
                '  "confidence_score": 0-100,\n'
                '  "missing_topics": ["..."],\n'
                '  "contradictions": ["..."],\n'
                '  "internet_queries": ["..."],\n'
                '  "document_queries": ["..."],\n'
                '  "fix_steps": ["..."]\n'
                '}'
            )
            validation_hint = 'Проверь аналитический текст на пробелы, противоречия и недоказанные выводы.'

        messages = [
            {
                'role': 'system',
                'content': (
                    'Ты validator для MTS Router. Проверь качество аналитики и верни результат строго в JSON.'
                ),
            },
            {
                'role': 'user',
                'content': (
                    f'Запрос:\n{prompt}\n\n'
                    f'Контекст для проверки:\n{context}\n\n'
                    f'{validation_hint}\n\n'
                    'Верни ТОЛЬКО JSON-объект без markdown и комментариев:\n'
                    f'{validation_schema}\n'
                    'Запросы должны быть конкретными и пригодными для web/document retrieval.'
                ),
            },
        ]
        temperature = 0
    else:
        messages = [
            {
                'role': 'system',
                'content': (
                    'Ты reviewer для MTS Router. Найди фактические ошибки, пробелы, противоречия и слабую структуру. '
                    'Отвечай компактно и прикладно.'
                ),
            },
            {
                'role': 'user',
                'content': (
                    f'Исходный запрос:\n{prompt}\n\n'
                    f'Черновик и контекст:\n{context}\n\n'
                    'Верни: 1) основные проблемы 2) конкретные правки 3) минимальный план улучшения.'
                ),
            },
        ]
        temperature = 0

    critique_text = await create_text_response(
        request=request,
        model=model_id,
        messages=messages,
        temperature=temperature,
    )
    return Artifact(
        id=operation.output_id,
        type='structured',
        source='operation',
        text=critique_text,
        metadata={
            'model_id': model_id,
            'stage': stage or None,
            'agent_role': agent_role or None,
        },
    )


async def _execute_generate_text_with_rounds(
    request: Request,
    form_data: dict[str, Any],
    operation,
    artifact_map: dict[str, Artifact],
    model_id: str,
    status_emitter: StatusEmitter | None,
    step: int,
    total_steps: int,
) -> Artifact:
    max_rounds = _to_positive_int(getattr(operation, 'max_rounds', 1), 1)
    stop_on_no_change = max(0, int(getattr(operation, 'stop_on_no_change', 0) or 0))

    previous_text = None
    no_change_streak = 0
    final_artifact = None

    for attempt in range(1, max_rounds + 1):
        if max_rounds > 1:
            await _emit_router_status(
                status_emitter,
                f'MTS Router: этап {step}/{total_steps} — раунд {attempt}/{max_rounds}',
                done=False,
                phase='operation',
                operation_id=operation.id,
                operation_kind=operation.kind,
                model_id=model_id,
                step=step,
                total_steps=total_steps,
                round=attempt,
                total_rounds=max_rounds,
            )

        artifact = await _execute_generate_text(
            request=request,
            form_data=form_data,
            operation=operation,
            artifact_map=artifact_map,
            model_id=model_id,
            previous_text=previous_text,
            attempt=attempt,
            max_rounds=max_rounds,
        )
        final_artifact = artifact
        current_text = artifact.text or ''

        if previous_text is not None:
            if _normalized_fingerprint(current_text) == _normalized_fingerprint(previous_text):
                no_change_streak += 1
            else:
                no_change_streak = 0

        previous_text = current_text
        if stop_on_no_change and no_change_streak >= stop_on_no_change:
            artifact.metadata['rounds_completed'] = attempt
            artifact.metadata['stopped_by_no_change'] = True
            artifact.metadata['no_change_streak'] = no_change_streak
            return artifact

    if final_artifact is None:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail='generate_text produced no output artifact.',
        )
    final_artifact.metadata['rounds_completed'] = max_rounds
    return final_artifact


async def build_execution_result(
    request: Request,
    form_data: dict[str, Any],
    user: UserModel | Any,
    policy_mode: str = 'balanced',
    dry_run: bool = False,
    status_emitter: StatusEmitter | None = None,
) -> ExecutionResult:
    if not request.app.state.MODELS:
        await get_all_models(request, user=user)

    registry = build_registry(request)
    if not registry.get('models'):
        await _emit_router_status(
            status_emitter,
            'MTS Router: не найдены доступные MWS-модели. Проверьте провайдер.',
            done=True,
            error=True,
            phase='plan',
        )
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail='MWS registry is empty. Make sure the MWS provider is configured and reachable.',
        )

    await _emit_router_status(
        status_emitter,
        'MTS Router: анализирую запрос и строю план',
        done=False,
        phase='plan',
        policy_mode=policy_mode,
    )

    execution_started_at = time.perf_counter()
    plan = build_plan(form_data, policy_mode=policy_mode)
    result = ExecutionResult(plan=plan)
    result.metrics = {
        'operations_total': len(plan.operations),
        'operations_completed': 0,
        'fallback_count': 0,
        'rounds_total': 0,
        'operation_latency_ms': {},
        'operation_models': {},
    }
    _add_trace(result, 'plan', 'completed', detail=f'Built deterministic plan: {plan.intent_mode}')
    await _emit_router_status(
        status_emitter,
        f'MTS Router: план готов ({plan.intent_mode}), этапов: {len(plan.operations)}',
        done=False,
        phase='plan',
        intent_mode=plan.intent_mode,
        operation_count=len(plan.operations),
    )

    if dry_run:
        await _emit_router_status(
            status_emitter,
            'MTS Router: dry-run завершён, выполнение этапов пропущено',
            done=True,
            phase='finalize',
            dry_run=True,
        )
        return result

    artifact_map = _artifact_map(result)
    total_operations = len(plan.operations)
    transient_unstable_models: set[str] = set()

    for operation_index, operation in enumerate(plan.operations, start=1):
        operation_started_at = time.perf_counter()
        assignment = select_operation_assignment(
            registry,
            operation.kind,
            plan.policy_mode,
            plan=plan,
            operation=operation,
        )
        result.selected_models[operation.id] = assignment
        candidate_models = [model_id for model_id in [assignment.primary, *assignment.fallbacks] if model_id]
        if transient_unstable_models and len(candidate_models) > 1:
            stable_models = [mid for mid in candidate_models if mid not in transient_unstable_models]
            unstable_models = [mid for mid in candidate_models if mid in transient_unstable_models]
            candidate_models = stable_models + unstable_models
        effective_primary_model = candidate_models[0] if candidate_models else assignment.primary

        if not candidate_models:
            await _emit_router_status(
                status_emitter,
                f'MTS Router: для этапа {operation.kind} нет доступных моделей',
                done=True,
                error=True,
                phase='operation',
                operation_id=operation.id,
                operation_kind=operation.kind,
            )
            raise HTTPException(
                status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
                detail=f'No MWS models available for operation {operation.kind}.',
            )

        _add_trace(
            result,
            'operation',
            'started',
            operation_id=operation.id,
            detail=f'Starting {operation.kind}',
            model_id=assignment.primary,
            step=operation_index,
            total_steps=total_operations,
        )
        await _emit_router_status(
            status_emitter,
            (
                f'MTS Router: этап {operation_index}/{total_operations} — {operation.kind} [{assignment.role}] '
                f'(основная модель: {effective_primary_model})'
            ),
            done=False,
            phase='operation',
            operation_id=operation.id,
            operation_kind=operation.kind,
            model_id=effective_primary_model,
            step=operation_index,
            total_steps=total_operations,
        )

        last_error = None
        output_artifact = None
        for model_id in candidate_models:
            try:
                if operation.kind == 'generate_text':
                    output_artifact = await _execute_generate_text_with_rounds(
                        request=request,
                        form_data=form_data,
                        operation=operation,
                        artifact_map=artifact_map,
                        model_id=model_id,
                        status_emitter=status_emitter,
                        step=operation_index,
                        total_steps=total_operations,
                    )
                elif operation.kind == 'analyze_image':
                    output_artifact = await _execute_analyze_image(request, operation, artifact_map, model_id)
                elif operation.kind == 'transcribe_audio':
                    output_artifact = await _execute_transcribe_audio(request, operation, artifact_map, model_id)
                elif operation.kind == 'retrieve_context':
                    output_artifact = await _execute_retrieve_context(
                        request,
                        form_data,
                        user,
                        operation,
                        artifact_map,
                        model_id,
                        status_emitter=status_emitter,
                        status_context={
                            'operation_id': operation.id,
                            'operation_kind': operation.kind,
                            'model_id': model_id,
                            'step': operation_index,
                            'total_steps': total_operations,
                        },
                    )
                elif operation.kind == 'critique_text':
                    output_artifact = await _execute_critique_text(request, operation, artifact_map, model_id)
                else:
                    raise HTTPException(
                        status_code=status.HTTP_400_BAD_REQUEST,
                        detail=f'Unsupported operation kind: {operation.kind}',
                    )

                result.selected_models[operation.id].selected = model_id
                break
            except Exception as exc:
                last_error = exc
                if _is_transient_provider_error(exc):
                    transient_unstable_models.add(model_id)
                result.metrics['fallback_count'] += 1
                _add_trace(
                    result,
                    'operation',
                    'fallback',
                    operation_id=operation.id,
                    model_id=model_id,
                    detail=str(exc),
                    step=operation_index,
                    total_steps=total_operations,
                )
                log.warning(
                    'MTS Router fallback: operation=%s model=%s err=%s',
                    operation.kind,
                    model_id,
                    _truncate_error(exc, max_length=400),
                )
                await _emit_router_status(
                    status_emitter,
                    f'MTS Router: резервная маршрутизация для этапа {operation.kind}',
                    done=False,
                    hidden=True,
                    phase='operation',
                    operation_id=operation.id,
                    operation_kind=operation.kind,
                    model_id=model_id,
                )

        if output_artifact is None:
            await _emit_router_status(
                status_emitter,
                (
                    f'MTS Router: этап {operation.kind} не выполнен. '
                    f'Последняя ошибка: {_truncate_error(last_error) if last_error else "unknown"}'
                ),
                done=True,
                error=True,
                phase='operation',
                operation_id=operation.id,
                operation_kind=operation.kind,
            )
            raise HTTPException(
                status_code=getattr(last_error, 'status_code', status.HTTP_500_INTERNAL_SERVER_ERROR),
                detail=str(last_error) if last_error else f'Operation failed: {operation.kind}',
            )

        result.output_artifacts.append(output_artifact)
        artifact_map[output_artifact.id] = output_artifact
        elapsed_ms = int((time.perf_counter() - operation_started_at) * 1000)
        rounds_completed = int((output_artifact.metadata or {}).get('rounds_completed') or 1)
        result.metrics['rounds_total'] += max(1, rounds_completed)
        result.metrics['operations_completed'] += 1
        result.metrics['operation_latency_ms'][operation.id] = elapsed_ms
        result.metrics['operation_models'][operation.id] = result.selected_models[operation.id].selected
        _add_trace(
            result,
            'operation',
            'completed',
            operation_id=operation.id,
            model_id=result.selected_models[operation.id].selected,
            detail=f'Completed {operation.kind}',
            step=operation_index,
            total_steps=total_operations,
            elapsed_ms=elapsed_ms,
            metadata={'rounds_completed': rounds_completed},
        )
        await _emit_router_status(
            status_emitter,
            (
                f'MTS Router: этап {operation_index}/{total_operations} завершён — {operation.kind} [{assignment.role}] '
                f'(модель: {result.selected_models[operation.id].selected})'
            ),
            done=False,
            phase='operation',
            operation_id=operation.id,
            operation_kind=operation.kind,
            model_id=result.selected_models[operation.id].selected,
            artifact_id=output_artifact.id,
            step=operation_index,
            total_steps=total_operations,
        )

    await _emit_router_status(
        status_emitter,
        'MTS Router: собираю финальный ответ',
        done=False,
        phase='finalize',
    )
    final_artifacts = [artifact_map[artifact_id] for artifact_id in plan.output_artifact_ids if artifact_id in artifact_map]
    final_text = next((artifact.text for artifact in reversed(final_artifacts) if artifact.text), None)
    result.final_output = FinalOutput(text=final_text, artifacts=final_artifacts)
    result.metrics['total_latency_ms'] = int((time.perf_counter() - execution_started_at) * 1000)
    _add_trace(result, 'finalize', 'completed', detail='Built final output payload.')
    await _emit_router_status(
        status_emitter,
        'MTS Router: выполнение завершено',
        done=True,
        phase='finalize',
    )

    return result


def build_openai_chat_completion_response(text: str, model: str = MTS_ROUTER_MODEL_ID) -> dict[str, Any]:
    created = int(time.time())
    return {
        'id': f'chatcmpl-{uuid.uuid4().hex[:16]}',
        'object': 'chat.completion',
        'created': created,
        'model': model,
        'choices': [
            {
                'index': 0,
                'message': {
                    'role': 'assistant',
                    'content': text,
                },
                'finish_reason': 'stop',
            }
        ],
        'usage': {
            'prompt_tokens': 0,
            'completion_tokens': 0,
            'total_tokens': 0,
        },
    }


def build_openai_streaming_response(text: str, model: str = MTS_ROUTER_MODEL_ID) -> StreamingResponse:
    created = int(time.time())
    chunk_id = f'chatcmpl-{uuid.uuid4().hex[:16]}'

    async def event_stream():
        first_chunk = {
            'id': chunk_id,
            'object': 'chat.completion.chunk',
            'created': created,
            'model': model,
            'choices': [
                {
                    'index': 0,
                    'delta': {
                        'role': 'assistant',
                        'content': text,
                    },
                    'finish_reason': None,
                }
            ],
        }
        yield f'data: {json.dumps(first_chunk, ensure_ascii=False)}\n\n'

        done_chunk = {
            'id': chunk_id,
            'object': 'chat.completion.chunk',
            'created': created,
            'model': model,
            'choices': [
                {
                    'index': 0,
                    'delta': {},
                    'finish_reason': 'stop',
                }
            ],
        }
        yield f'data: {json.dumps(done_chunk, ensure_ascii=False)}\n\n'
        yield 'data: [DONE]\n\n'

    return StreamingResponse(event_stream(), media_type='text/event-stream')


async def execute_orchestrated_chat_completion(
    request: Request,
    form_data: dict[str, Any],
    user: UserModel | Any,
):
    policy_mode = (
        (form_data.get('params') or {}).get('orchestration_mode')
        or getattr(request.app.state, 'MTS_ROUTER_DEFAULT_POLICY_MODE', 'balanced')
    )

    metadata = form_data.get('metadata') or getattr(request.state, 'metadata', {})
    params = dict(form_data.get('params') or {})
    show_router_status = params.get('show_router_status')
    emit_public_router_status = True if show_router_status is None else _to_bool(show_router_status)

    is_background_task = isinstance(metadata, dict) and bool(metadata.get('task'))
    event_emitter = None if is_background_task else (get_event_emitter(metadata) if isinstance(metadata, dict) else None)
    status_emitter = event_emitter if emit_public_router_status else None

    try:
        result = await build_execution_result(
            request=request,
            form_data=form_data,
            user=user,
            policy_mode=policy_mode,
            dry_run=False,
            status_emitter=status_emitter,
        )
    except Exception as exc:
        await _emit_router_status(
            status_emitter,
            f'MTS Router: выполнение завершилось ошибкой: {_truncate_error(exc)}',
            done=True,
            error=True,
            phase='finalize',
        )
        raise

    text = result.final_output.text or 'MTS Router did not produce a textual response.'

    plan_metadata = dict(getattr(result.plan, 'metadata', {}) or {})
    should_try_attach_presentation = _to_bool(plan_metadata.get('presentation_generation'))
    deck: list[dict[str, Any]] = []

    if should_try_attach_presentation:
        deck = _extract_presentation_deck(text)
    else:
        # Fallback: if planner metadata was not set but the model still returned a slide deck JSON.
        deck = _extract_presentation_deck(text)
        if deck and len(deck) >= 3:
            lowered = text.lower()
            should_try_attach_presentation = any(
                marker in lowered
                for marker in (
                    'deck json',
                    'финальный deck',
                    'presentation',
                    'презентац',
                    'слайд',
                    'slides',
                )
            )

    if should_try_attach_presentation and len(deck) < 3:
        recovered_deck = _extract_presentation_deck_from_artifacts(
            result.output_artifacts,
            preferred_ids=[
                'presentation_content_2',
                'presentation_content_1',
                'out_text_1',
            ],
            min_slides=3,
        )
        if recovered_deck:
            deck = recovered_deck
            await _emit_router_status(
                status_emitter,
                'MTS Router: восстановил deck из промежуточных артефактов и продолжаю сборку презентации',
                done=False,
                hidden=True,
                phase='finalize',
                action='presentation_deck_recovered',
                slides=len(deck),
            )

    if should_try_attach_presentation and deck and _is_presentation_refusal_text(text):
        text = _build_presentation_recovery_text(deck)

    if should_try_attach_presentation and deck:
        try:
            render_deck = _normalize_deck_for_render(deck)
            if not render_deck:
                render_deck = deck

            pptx_bytes, renderer_name = await asyncio.to_thread(_render_presentation_pptx, render_deck)
            attachment = await _attach_presentation_file_to_chat(
                request=request,
                user=user,
                metadata=metadata if isinstance(metadata, dict) else {},
                deck=render_deck,
                pptx_bytes=pptx_bytes,
                renderer=renderer_name,
                event_emitter=event_emitter,
            )
            if attachment:
                download_url = str(attachment.get('download_url') or '').strip()
                download_line = (
                    f'\nСкачать: [{attachment.get("name")}]({download_url})'
                    if download_url
                    else ''
                )
                text = (
                    f'{text.rstrip()}\n\n'
                    'Презентация `.pptx` прикреплена к этому сообщению '
                    f'(откройте файл под ответом).{download_line}'
                )
        except Exception as exc:
            log.warning('Failed to build/attach presentation file: %s', _truncate_error(exc, max_length=400))

    if form_data.get('stream'):
        return build_openai_streaming_response(text=text, model=MTS_ROUTER_MODEL_ID)
    return build_openai_chat_completion_response(text=text, model=MTS_ROUTER_MODEL_ID)
